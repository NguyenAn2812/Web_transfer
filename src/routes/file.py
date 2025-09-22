import os, html, csv, io, re

from datetime import datetime
from flask import Blueprint, render_template, request, redirect, url_for, current_app, send_file, abort
from ..utils.db import get_db
from ..utils.device import get_device_name, get_all_devices
from ..utils.file import allowed_file, get_file_type, content_hash, safe_filename, shorten_name, ext_of

bp = Blueprint('file', __name__)

@bp.route('/upload', methods=['GET', 'POST'])
def upload():
    sid = request.cookies.get(current_app.config['SESSION_COOKIE_NAME'])
    device_name = get_device_name(sid) if sid else None
    devices = get_all_devices(exclude_id=sid)

    if request.method == 'POST':
        file = request.files.get('file')
        mode = (request.form.get('mode') or '').strip()         # 'public' | 'private'
        receiver = (request.form.get('receiver') or '').strip() # session_id người nhận (private)
        display_name = (request.form.get('display_name') or '').strip()
        title = (request.form.get('title') or '').strip()
        folder = (request.form.get('folder') or 'root').strip()
        # sanitize folder: chỉ chữ, số, _ - / (cấp con), chống path traversal
        folder = re.sub(r'[^A-Za-z0-9_\-\/]', '', folder) or 'root'
        folder = folder.strip('/').strip('\\') or 'root'
        if not file or not mode:
            return render_template('upload.html', has_session=bool(sid), device_name=device_name,
                                   devices=devices, session_id=sid, error='Thiếu thông tin')

        if not allowed_file(file.filename):
            return render_template('upload.html', has_session=bool(sid), device_name=device_name,
                                   devices=devices, session_id=sid, error='Tên file không hợp lệ')

        # sanitize + rút gọn hiển thị
        orig_name = safe_filename(file.filename)
        if not display_name: display_name = orig_name
        if not title: title = orig_name
        display_name = shorten_name(display_name)
        title = shorten_name(title, limit=10)

        # hash nội dung để kiểm tra trùng
        fhash, fsize = content_hash(file)

        db = get_db()
        folders = [r['folder'] for r in db.execute(
            "SELECT DISTINCT folder FROM files ORDER BY folder ASC").fetchall()]
        return render_template('upload.html', has_session=bool(sid), device_name=device_name,
                            devices=devices, session_id=sid, folders=folders)
        cur = db.execute('SELECT id, file_path FROM files WHERE file_hash=?', (fhash,))
        dup = cur.fetchone()
        if dup:
            # file đã tồn tại → chỉ thêm metadata mới (tham chiếu chung 1 path)
            db.execute('''INSERT INTO files (filename, display_name, title, uploader_id, uploader_name,
                    is_public, receiver_id, upload_time, file_path, file_type, file_hash, file_size, folder)
                    SELECT ?,?,?,?,?,?,?,?, file_path, ?, file_hash, ?, ? FROM files WHERE id=?''',
                    (orig_name, display_name, title, sid, device_name,
                    1 if mode=='public' else 0, receiver if mode=='private' and receiver else None,
                    datetime.now().isoformat(), get_file_type(orig_name), fsize, folder, dup['id']))

            db.commit()
            return redirect(url_for('file.files'))

        # chưa có → lưu file theo hash
        upload_dir = os.path.join(current_app.config['UPLOAD_FOLDER'], folder)  # NEW: subfolder

        ext = os.path.splitext(orig_name)[1].lower()
        hash_name = f"{fhash}{ext if ext else ''}"
        file_path = os.path.join(upload_dir, hash_name)
        os.makedirs(upload_dir, exist_ok=True)
        file.save(file_path)

        db.execute('''INSERT INTO files (filename, display_name, title, uploader_id, uploader_name,
            is_public, receiver_id, upload_time, file_path, file_type, file_hash, file_size, folder)
            VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?)''',
           (orig_name, display_name, title, sid, device_name,
            1 if mode=='public' else 0, receiver if mode=='private' and receiver else None,
            datetime.now().isoformat(), file_path, get_file_type(orig_name), fhash, fsize, folder))

        db.commit()
        return redirect(url_for('file.files'))

    return render_template('upload.html', has_session=bool(sid), device_name=device_name,
                           devices=devices, session_id=sid)

@bp.get('/files')
def files():
    sid = request.cookies.get(current_app.config['SESSION_COOKIE_NAME'])
    device_name = get_device_name(sid) if sid else None
    folder = (request.args.get('folder') or '').strip() or None

    db = get_db()
    if folder:
        cur = db.execute('''SELECT * FROM files
                            WHERE (is_public=1 OR receiver_id=?)
                              AND folder=?
                            ORDER BY id DESC''', (sid, folder))
    else:
        cur = db.execute('''SELECT * FROM files
                            WHERE is_public=1 OR receiver_id=?
                            ORDER BY id DESC''', (sid,))
    rows = cur.fetchall()

    # danh sách folder để hiển thị filter
    folders = [r['folder'] for r in db.execute(
        "SELECT DISTINCT folder FROM files ORDER BY folder ASC").fetchall()]

    return render_template('files.html', has_session=bool(sid), device_name=device_name,
                           files=rows, session_id=sid, folder=folder, folders=folders)

@bp.get('/download/<int:file_id>')
def download(file_id):
    # SỬA: truy vấn theo id thật trong DB, không dùng index danh sách tạm
    db = get_db()
    row = db.execute('SELECT * FROM files WHERE id=?', (file_id,)).fetchone()
    if not row:
        return abort(404)

    # kiểm tra quyền xem: public hoặc là người nhận
    sid = request.cookies.get(current_app.config['SESSION_COOKIE_NAME'])
    if not (row['is_public'] == 1 or (sid and row['receiver_id'] == sid) or (sid and row['uploader_id'] == sid)):
        return abort(403)

    fpath = row['file_path']

    # Fallback: nếu path cũ không tồn tại, thử ghép lại từ UPLOAD_FOLDER + basename
    if not fpath or not os.path.exists(fpath):
        basename = os.path.basename(row['file_path'] or row['filename'] or '')
        sub = row['folder'] or 'root'
        alt = os.path.join(current_app.config['UPLOAD_FOLDER'], sub, basename) if basename else None

        if alt and os.path.exists(alt):
            fpath = alt

    if not fpath or not os.path.exists(fpath):
        return abort(404)

    return send_file(fpath, as_attachment=True, download_name=row['filename'])


@bp.get('/delete/<int:file_id>')
def delete_file(file_id):
    sid = request.cookies.get(current_app.config['SESSION_COOKIE_NAME'])
    db = get_db()
    row = db.execute('SELECT * FROM files WHERE id=?', (file_id,)).fetchone()
    if not row:
        return abort(404)
    if row['uploader_id'] != sid:
        return abort(403)

    # xóa record hiện tại
    db.execute('DELETE FROM files WHERE id=?', (file_id,))
    db.commit()

    # nếu không còn record nào tham chiếu cùng hash thì xóa file vật lý
    left = db.execute('SELECT COUNT(*) AS c FROM files WHERE file_hash=?', (row['file_hash'],)).fetchone()['c']
    if left == 0 and row['file_path'] and os.path.exists(row['file_path']):
        try:
            os.remove(row['file_path'])
        except OSError:
            pass

    return redirect(url_for('file.files'))


@bp.get('/preview/<int:file_id>')
def preview_file(file_id):
    db = get_db()
    row = db.execute('SELECT * FROM files WHERE id=?', (file_id,)).fetchone()
    if not row:
        return abort(404)

    # quyền xem
    sid = request.cookies.get(current_app.config['SESSION_COOKIE_NAME'])
    if not (row['is_public'] == 1 or (sid and row['receiver_id'] == sid) or (sid and row['uploader_id'] == sid)):
        return abort(403)

    ftype = row['file_type']
    fpath = row['file_path']

    # Fallback: nếu path cũ không tồn tại, thử ghép lại từ UPLOAD_FOLDER + basename
    if not fpath or not os.path.exists(fpath):
        basename = os.path.basename(row['file_path'] or row['filename'] or '')
        sub = row['folder'] or 'root'
        alt = os.path.join(current_app.config['UPLOAD_FOLDER'], sub, basename) if basename else None

        if alt and os.path.exists(alt):
            fpath = alt

    if not fpath or not os.path.exists(fpath):
        return abort(404)

    # đường dẫn public qua /uploads/<filename>
    public_url = url_for('uploaded_file', filename=os.path.basename(fpath))


    # ---------- IMAGE ----------
    if ftype == 'image':
        return f'''
        <html><head><meta charset="utf-8"><title>Preview</title></head>
        <body style="margin:0;background:#111;color:#eee">
          <img src="{public_url}" style="max-width:100vw;max-height:100vh;object-fit:contain;display:block;margin:auto">
        </body></html>'''

    # ---------- PDF ----------
    if ftype == 'pdf':
        return f'''
        <html><head><meta charset="utf-8"><title>Preview</title></head>
        <body style="margin:0">
          <embed src="{public_url}" type="application/pdf" width="100%" height="100%" style="height:100vh;border:none"/>
        </body></html>'''

    # ---------- TEXT / CSV ----------
    if ftype == 'text':
        ext = ext_of(row['filename'])
        if ext == 'csv':
            # Hiển thị 200 dòng đầu thành bảng
            try:
                with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = []
                    for i, r in enumerate(reader):
                        rows.append(r[:50])  # tối đa 50 cột
                        if i >= 200: break   # tối đa 200 dòng
                # dựng bảng HTML đơn giản
                def td(x): return f"<td style='border:1px solid #ddd;padding:4px'>{html.escape(str(x))}</td>"
                trs = []
                for i, r in enumerate(rows):
                    cells = ''.join(td(x) for x in r)
                    trs.append(f"<tr>{cells}</tr>")
                table_html = "<table style='border-collapse:collapse;font-family:monospace;font-size:13px'>" + ''.join(trs) + "</table>"
                return f"<html><head><meta charset='utf-8'><title>CSV Preview</title></head><body>{table_html}</body></html>"
            except Exception as e:
                return f"<pre>Lỗi đọc CSV: {html.escape(str(e))}</pre>", 500
        else:
            # text thường: đọc 200 KB đầu
            try:
                with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(200_000)
                return f"<pre style='white-space:pre-wrap'>{html.escape(content)}</pre>"
            except Exception as e:
                return f"<pre>Lỗi đọc text: {html.escape(str(e))}</pre>", 500

    # ---------- XLS/XLSX (nếu có openpyxl) ----------
    if ftype == 'xls':
        ext = ext_of(row['filename'])
        if ext in {'xlsx','xls'}:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(fpath, read_only=True, data_only=True)
                ws = wb.active
                max_r, max_c = 100, 30  # giới hạn preview
                trs = []
                for r in ws.iter_rows(min_row=1, max_row=max_r, min_col=1, max_col=max_c, values_only=True):
                    cells = ''.join(f"<td style='border:1px solid #ddd;padding:4px'>{html.escape('' if v is None else str(v))}</td>" for v in r)
                    trs.append(f"<tr>{cells}</tr>")
                table_html = "<table style='border-collapse:collapse;font-family:monospace;font-size:13px'>" + ''.join(trs) + "</table>"
                return f"<html><head><meta charset='utf-8'><title>Excel Preview</title></head><body>{table_html}</body></html>"
            except ImportError:
                return "<p>Chưa cài <code>openpyxl</code>. Cài thêm để xem trước Excel: <code>pip install openpyxl</code></p>"
            except Exception as e:
                return f"<pre>Lỗi đọc Excel: {html.escape(str(e))}</pre>", 500

    # ---------- PPT/PPTX (nếu có python-pptx) ----------
    if ftype == 'ppt':
        try:
            from pptx import Presentation
            prs = Presentation(fpath)
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                block = "<h3>Slide {}</h3><div style='border:1px solid #ddd;padding:6px;margin-bottom:8px'>{}</div>".format(
                    i, "<br>".join(html.escape(t) for t in texts[:20]))
                parts.append(block)
                if i >= 20:  # giới hạn 20 slide đầu
                    break
            if not parts:
                parts.append("<p>(Không tìm thấy text — file có thể toàn hình/đồ thị)</p>")
            return "<html><head><meta charset='utf-8'><title>PPT Preview</title></head><body>{}</body></html>".format("".join(parts))
        except ImportError:
            return "<p>Chưa cài <code>python-pptx</code>. Cài thêm để xem trước PPT: <code>pip install python-pptx</code></p>"
        except Exception as e:
            return f"<pre>Lỗi đọc PPT: {html.escape(str(e))}</pre>", 500

    # ---------- OTHER ----------
    return f"""
    <html><head><meta charset="utf-8"><title>Preview</title></head>
    <body>
      <p>Không hỗ trợ xem trước định dạng này. Bạn có thể <a href="{url_for('file.download', file_id=file_id)}">tải xuống</a>.</p>
    </body></html>"""
